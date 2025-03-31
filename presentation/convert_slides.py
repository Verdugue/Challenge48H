"""
Script pour convertir les slides Markdown en présentation PowerPoint.
Utilise python-pptx pour la conversion.
"""

from pptx import Presentation
import re

def create_title_slide(prs, title, subtitle=None):
    """Crée une slide de titre."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle

def create_content_slide(prs, title, content):
    """Crée une slide de contenu."""
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    content_shape.text = content

def convert_markdown_to_pptx(markdown_file, output_file):
    """Convertit un fichier Markdown en présentation PowerPoint."""
    prs = Presentation()
    
    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Diviser le contenu en slides
    slides = content.split('---')
    
    for slide in slides:
        if not slide.strip():
            continue
            
        # Extraire le titre et le contenu
        lines = slide.strip().split('\n')
        title = lines[0].replace('# ', '')
        content = '\n'.join(lines[1:]).strip()
        
        # Créer la slide
        if '##' in content:
            # Slide avec sous-titre
            create_title_slide(prs, title, content.split('##')[1].strip())
        else:
            # Slide de contenu
            create_content_slide(prs, title, content)
    
    # Sauvegarder la présentation
    prs.save(output_file)

if __name__ == '__main__':
    convert_markdown_to_pptx('presentation/slides.md', 'presentation/presentation.pptx') 